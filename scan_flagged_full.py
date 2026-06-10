"""Get full text of flagged shapes."""
from pptx import Presentation

prs = Presentation(r'D:\Hospital Resource Optimization with AI Final Hro.pptx')

targets = {
    2: ['Text 3', 'Text 4'],           # S03
    8: ['Text 2', 'Text 6', 'Text 14'], # S09
    14: ['Text 1'],                     # S15
    15: ['Text 16', 'Text 25', 'Text 29', 'Text 31'],  # S16
}

for slide_idx, shape_names in targets.items():
    slide = prs.slides[slide_idx]
    print(f'=== S{slide_idx+1} ===')
    for shape in slide.shapes:
        if shape.name in shape_names:
            print(f'  [{shape.name}]: {repr(shape.text_frame.text)}')
    print()
