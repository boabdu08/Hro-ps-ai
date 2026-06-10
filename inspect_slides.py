"""Inspect slide structure and table contents in detail."""
from pptx import Presentation

prs = Presentation(r'D:\Hospital Resource Optimization with AI Final Hro.pptx')

# Check S13 (index 12) for tables
for target_idx in [8, 9, 12, 13, 14, 19, 20]:
    slide = prs.slides[target_idx]
    print(f'=== S{target_idx+1} Shapes ===')
    for j, shape in enumerate(slide.shapes):
        has_t = shape.has_text_frame
        has_tbl = shape.has_table
        print(f'  Shape {j}: name={shape.name!r}, has_text={has_t}, has_table={has_tbl}')
        if has_t:
            txt = shape.text_frame.text.strip()[:80]
            print(f'    Text: {txt.encode("ascii","replace").decode()}')
        if has_tbl:
            print(f'    Table rows: {len(shape.table.rows)}, cols: {len(shape.table.columns)}')
            for ri, row in enumerate(shape.table.rows):
                cells = [cell.text_frame.text.strip()[:30].encode("ascii","replace").decode() for cell in row.cells]
                print(f'    Row {ri}: {cells}')
        if hasattr(shape, 'notes_slide') and shape.has_notes_slide:
            pass
    # Check notes
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            print(f'  NOTES: {notes[:100].encode("ascii","replace").decode()}')
    print()
