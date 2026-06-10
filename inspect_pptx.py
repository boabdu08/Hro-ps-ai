"""Quick structural inspection of PPTX source files."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

def inspect(path: str) -> None:
    prs = Presentation(path)
    print(f"\n=== {Path(path).name} ===")
    print(f"Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t[:80])
        if texts:
            line = " | ".join(texts[:4])
            print(f"  S{i:02d}: {line}".encode("ascii", "replace").decode("ascii"))

for p in [r"D:\poster msa.pptx",
          r"D:\Hospital Resource Optimization with AI Final Hro.pptx"]:
    try:
        inspect(p)
    except Exception as e:
        print(f"Error reading {p}: {e}")
