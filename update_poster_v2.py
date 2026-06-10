"""
Phase 10 — Comprehensive poster rebuild.

Fixes:
- Global text subs (dataset, optimizer, weights)
- KPI cards: replace fabricated surgery/OR metrics with real model metrics
- Table 83 (Forecast MAPE): update 72-h column to canonical values; note other columns
- Table 85 (OR Scheduling): replace with canonical evaluation metrics table
- Insert verified figures: architecture, forecasting pipeline, results metrics, screenshots
"""
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import re

BASE  = Path(r"D:\Hro new dashboard")
FIGS  = BASE / "thesis_figures"
SHOTS = FIGS / "screenshots"

SRC = r"D:\poster msa.pptx"
DST = str(BASE / "HRO-PS_Poster_FINAL_v2.pptx")

# -----------------------------------------------------------------------
# GLOBAL TEXT SUBS (applied to every run in every shape)
# -----------------------------------------------------------------------
SUBS = [
    ("29 302",   "17,520"),
    ("29,302",   "17,520"),
    ("29302",    "17,520"),
    ("3+ years", "2 years"),
    ("3 years",  "2 years"),
    ("168-hour", "24-hour"),
    ("0.6/0.4",  "0.80/0.20"),
    ("0.7/0.3",  "0.80/0.20"),
    ("rule-based optimization", "MILP optimization (scipy.optimize.milp)"),
    ("rule-based resource",     "MILP-based resource"),
    ("OR-Tools",                "scipy.optimize.milp"),
    ("TimescaleDB",             "PostgreSQL"),
    ("React frontend",          "Streamlit dashboard"),
    ("108 tests",               "128 tests"),
    ("108 unit tests",          "128 unit tests"),
    ("25.5% improvement",       "~51% lower MAE vs ARIMAX"),
    # Hybrid weight corrections
    ("0.6/0.4", "0.80/0.20"),
    ("60%/40%", "80%/20%"),
]

REGEX_SUBS = [
    (re.compile(r'Hybrid\s+(is|model is)\s+(the\s+)?most accurate', re.I),
     "LSTM is the most accurate model; Hybrid is deployed for robustness"),
    (re.compile(r'\b94\.17\s*%'), "5.52% MAPE"),
]

# -----------------------------------------------------------------------
# KPI CARD REPLACEMENTS (shape-index → new_value, new_label)
# These are the big-number summary cards with fabricated surgery/OR stats
# -----------------------------------------------------------------------
KPI_REPLACEMENTS = {
    # Top-area KPI cards
    16: "6.07%",
    17: "Hybrid MAPE at 72-h horizon (deployed for robustness)",
    20: "7.65 pts",
    21: "LSTM MAE — most accurate model on test set",
    24: "< 5 sec",
    25: "MILP solver time (scipy.optimize.milp, 5-dept allocation)",
    # Results-section KPI cards
    69: "5.52%",
    70: "LSTM MAPE — best model, 72-h test set",
    72: "8.31 pts",
    73: "Hybrid MAE — deployed operational forecaster",
    75: "10.22 pts",
    76: "Hybrid RMSE — 0.80 LSTM + 0.20 ARIMAX blend",
}

# -----------------------------------------------------------------------
# TABLE FIXES
# -----------------------------------------------------------------------

def fix_table_83(shape):
    """Update Table 1 (Forecast MAPE) 72-h column to canonical test values."""
    tbl = shape.table
    # Expected structure:
    #   Row 0: headers (Model, 24-h, 48-h, 72-h)
    #   Row 1: Moving-Avg (baseline)
    #   Row 2: ARIMAX    → 72-h: 12.33%  (canonical test MAPE)
    #   Row 3: LSTM      → 72-h: 5.52%   (canonical test MAPE)
    #   Row 4: Hybrid    → 72-h: 6.07%   (canonical test MAPE)
    canonical_72h = {
        "ARIMAX":       "12.33%",
        "LSTM":         "5.52%",
        "Hybrid (ours)":"6.07%",
    }
    for row in tbl.rows:
        model_cell = row.cells[0].text_frame.text.strip()
        if model_cell in canonical_72h:
            cell = row.cells[3]  # 72-h column
            tf = cell.text_frame
            updated = False
            for para in tf.paragraphs:
                for run in para.runs:
                    if not updated:
                        run.text = canonical_72h[model_cell]
                        updated = True
                    else:
                        run.text = ""
    print("  Table 83 (Forecast MAPE): 72-h column updated to canonical test metrics")


def fix_table_85(shape):
    """Replace Table 2 (fabricated OR scheduling) with canonical evaluation metrics."""
    tbl = shape.table
    # New content: Model Evaluation Metrics
    new_data = [
        ["Model",      "MAE (patients)", "RMSE",  "MAPE",   "Role"],
        ["LSTM",       "7.65",           "9.58",  "5.52%",  "Most accurate"],
        ["ARIMAX",     "15.63",          "19.33", "12.33%", "Baseline"],
        ["Hybrid 0.80/0.20", "8.31",     "10.22", "6.07%",  "Deployed"],
    ]
    for ri, row in enumerate(tbl.rows):
        if ri >= len(new_data):
            break
        for ci, cell in enumerate(row.cells):
            if ci >= len(new_data[ri]):
                break
            cell_tf = cell.text_frame
            # Set text in first paragraph first run
            for para in cell_tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            if cell_tf.paragraphs:
                para = cell_tf.paragraphs[0]
                if para.runs:
                    para.runs[0].text = new_data[ri][ci]
                else:
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    # Add a run
                    r_el = etree.SubElement(para._p, qn("a:r"))
                    rPr = etree.SubElement(r_el, qn("a:rPr"), attrib={"lang": "en-US"})
                    t_el = etree.SubElement(r_el, qn("a:t"))
                    t_el.text = new_data[ri][ci]
    # Update label above table
    print("  Table 85 (OR Scheduling → Model Evaluation): updated to canonical MAE/RMSE/MAPE")


# -----------------------------------------------------------------------
# HELPERS
# -----------------------------------------------------------------------

def patch_run(run) -> bool:
    original = run.text
    text = original
    for old, new in SUBS:
        text = text.replace(old, new)
    for pattern, new in REGEX_SUBS:
        text = pattern.sub(new, text)
    if text != original:
        run.text = text
        return True
    return False


def patch_text_frame(tf) -> int:
    changed = 0
    for para in tf.paragraphs:
        for run in para.runs:
            if patch_run(run):
                changed += 1
    return changed


def set_shape_text(slide, shape_idx, new_text):
    from pptx.oxml.ns import qn
    from lxml import etree
    shape = list(slide.shapes)[shape_idx]
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Clear all existing runs
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    # Set text in first paragraph
    if tf.paragraphs:
        para = tf.paragraphs[0]
        if para.runs:
            para.runs[0].text = new_text
        else:
            # No runs exist — add one via lxml
            r_el = etree.SubElement(para._p, qn("a:r"))
            rPr = etree.SubElement(r_el, qn("a:rPr"), attrib={"lang": "en-US"})
            t_el = etree.SubElement(r_el, qn("a:t"))
            t_el.text = new_text
    safe_name = shape.name.encode('ascii', 'replace').decode()
    safe_text = new_text.encode('ascii', 'replace').decode()
    print(f"  Shape {shape_idx} [{safe_name}]: -> {safe_text!r}")


def add_image(slide, img_path, left_in, top_in, width_in, height_in=None):
    from pptx.util import Inches
    left  = Inches(left_in)
    top   = Inches(top_in)
    width = Inches(width_in)
    if height_in:
        slide.shapes.add_picture(img_path, left, top, width, Inches(height_in))
    else:
        slide.shapes.add_picture(img_path, left, top, width)


# -----------------------------------------------------------------------
# MAIN
# -----------------------------------------------------------------------

def main():
    prs = Presentation(SRC)
    slide = prs.slides[0]

    sw_in = prs.slide_width  / 914400
    sh_in = prs.slide_height / 914400
    print(f"Poster: {sw_in:.1f} x {sh_in:.1f} inches, {len(list(slide.shapes))} shapes")

    # 1. Global text patches
    print("\n--- Phase 1: Global text patches ---")
    total = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            total += patch_text_frame(shape.text_frame)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    total += patch_text_frame(cell.text_frame)
    print(f"  {total} runs changed globally")

    # 2. KPI card replacements
    print("\n--- Phase 2: KPI card replacements ---")
    shapes = list(slide.shapes)
    for shape_idx, new_text in KPI_REPLACEMENTS.items():
        if shape_idx < len(shapes):
            set_shape_text(slide, shape_idx, new_text)

    # 3. Fix tables
    print("\n--- Phase 3: Table fixes ---")
    for shape in slide.shapes:
        if shape.has_table:
            if shape.name == "Table 83":
                fix_table_83(shape)
            elif shape.name == "Table 85":
                fix_table_85(shape)

    # 4. Fix poster label caption for Table 85
    # Shape 83 = "TextBox 84" = "Table 2 — OR Scheduling — HRO-PS vs. Manual"
    try:
        set_shape_text(slide, 83,
            "Table 2  —  Model Evaluation Metrics (test set, chronological split)")
    except Exception as e:
        print(f"  Table 2 label: {e}")

    # 5. Insert verified figures
    print("\n--- Phase 4: Inserting verified figures ---")

    # Replace/supplement architecture area (Section 05 — roughly left column, middle)
    arch_png = str(FIGS / "fig_architecture.png")
    if Path(arch_png).exists():
        try:
            add_image(slide, arch_png,
                      left_in=0.3, top_in=6.2, width_in=4.0, height_in=2.8)
            print("  Architecture PNG added (left column, section 05)")
        except Exception as e:
            print(f"  Architecture PNG error: {e}")

    # Forecasting pipeline (middle column)
    fp_png = str(FIGS / "fig_forecasting_pipeline.png")
    if Path(fp_png).exists():
        try:
            add_image(slide, fp_png,
                      left_in=4.5, top_in=6.2, width_in=4.2, height_in=2.8)
            print("  Forecasting pipeline PNG added (middle column)")
        except Exception as e:
            print(f"  Forecasting pipeline PNG error: {e}")

    # Results metrics (right of results section)
    met_png = str(FIGS / "fig_results_metrics.png")
    if Path(met_png).exists():
        try:
            add_image(slide, met_png,
                      left_in=4.5, top_in=9.3, width_in=4.2, height_in=2.2)
            print("  Results metrics PNG added")
        except Exception as e:
            print(f"  Results metrics PNG error: {e}")

    # Command Center screenshot (section 07/08 — right column, demo section)
    cmd_jpg = str(SHOTS / "01_command_center.jpg")
    if Path(cmd_jpg).exists():
        try:
            add_image(slide, cmd_jpg,
                      left_in=9.0, top_in=6.5, width_in=4.0, height_in=2.5)
            print("  Command Center screenshot added")
        except Exception as e:
            print(f"  Command Center screenshot error: {e}")

    # Optimization screenshot (section 07 — demo)
    opt_jpg = str(SHOTS / "05_optimization.jpg")
    if Path(opt_jpg).exists():
        try:
            add_image(slide, opt_jpg,
                      left_in=9.0, top_in=9.3, width_in=4.0, height_in=2.5)
            print("  Optimization screenshot added")
        except Exception as e:
            print(f"  Optimization screenshot error: {e}")

    # Evaluation screenshot (section 06 — results) — shows the canonical metrics table live
    eval_jpg = str(SHOTS / "06_evaluation.jpg")
    if Path(eval_jpg).exists():
        try:
            add_image(slide, eval_jpg,
                      left_in=0.3, top_in=9.3, width_in=4.0, height_in=2.5)
            print("  Evaluation screenshot added (shows canonical metrics table)")
        except Exception as e:
            print(f"  Evaluation screenshot error: {e}")

    # 5. Save
    prs.save(DST)
    sz = Path(DST).stat().st_size // 1024
    print(f"\nSaved: {DST}  ({sz} KB)")
    print("Done.")


if __name__ == "__main__":
    main()
