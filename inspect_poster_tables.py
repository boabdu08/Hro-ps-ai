"""Inspect poster tables and key KPI text boxes."""
from pptx import Presentation

prs = Presentation(r"D:\poster msa.pptx")
slide = prs.slides[0]

# Tables
print("=== TABLES ===")
for shape in slide.shapes:
    if shape.has_table:
        print(f"\nTable: {shape.name}")
        for ri, row in enumerate(shape.table.rows):
            cells = [cell.text_frame.text.strip()[:30].encode("ascii","replace").decode()
                     for cell in row.cells]
            print(f"  Row {ri}: {cells}")

# Key text shapes
print("\n=== KEY TEXT BOXES ===")
kpi_indices = [5,6,16,17,18,20,21,22,24,25,26,28,29,35,41,47,53,54,61,63,69,70,71,72,73,74,75,76,77,91,97]
for i, shape in enumerate(slide.shapes):
    if shape.has_text_frame and i in kpi_indices:
        t = shape.text_frame.text.strip()
        if t:
            print(f"  Shape {i:3d} [{shape.name}]: {repr(t[:100])}")
