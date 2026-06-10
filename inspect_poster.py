"""Inspect poster slide structure."""
from pptx import Presentation

prs = Presentation(r"D:\poster msa.pptx")
print(f"Slides: {len(prs.slides)}")
print(f"Slide size: {prs.slide_width/914400:.1f} x {prs.slide_height/914400:.1f} inches")

slide = prs.slides[0]
print(f"\nShapes: {len(list(slide.shapes))}")
for i, shape in enumerate(slide.shapes):
    has_t = shape.has_text_frame
    has_tbl = shape.has_table
    has_img = not has_t and not has_tbl
    t = ""
    if has_t:
        t = shape.text_frame.text.strip()[:80].encode("ascii","replace").decode()
    print(f"  {i:3d}  {shape.name!r:20s}  text={has_t}  tbl={has_tbl}  img={has_img}  |  {t}")
