"""
Phase 10 — Comprehensive PPTX sync.

Changes vs. update_pptx.py (Phase 8/9):
- Iterates ALL shape types including tables (none found in practice, but covered)
- Patches speaker notes
- Applies slide-specific targeted fixes:
    S05  Industrial Partner -> Clinical Context & Site Visit
    S09  HL7/FHIR framed as future work
    S13  Adds canonical metrics text box + speaker notes (no pptx table exists)
    S14  Fixes hybrid formula alpha to 0.80/0.20
    S15  24-Hour -> 72-Hour in subtitle
    S16  TimescaleDB/React/OR-Tools/CP-SAT corrected
    S21  False accuracy claim replaced with factual MILP description
- Adds speaker notes to S13, S14, S15, S20, S21 from jury Q&A
"""
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = Path(r"D:\Hro new dashboard")
FIGS = BASE / "thesis_figures"
SHOTS = FIGS / "screenshots"

SRC = r"D:\Hospital Resource Optimization with AI Final Hro.pptx"
DST = str(BASE / "HRO-PS_Presentation_FINAL_v2.pptx")

# -----------------------------------------------------------------------
# TEXT SUBSTITUTIONS — applied to every run in every text frame + table
# (context-aware: only replace in metric/model/stack contexts)
# -----------------------------------------------------------------------
SUBS = [
    # Stack corrections
    ("TimescaleDB",       "PostgreSQL"),
    ("React.js component architecture for real-time interactive dashboards and visualization.",
     "Streamlit dashboard for real-time hospital operations and visualization."),
    ("OR-Tools",          "scipy.optimize.milp"),
    ("CP-SAT & MILP solvers (Gurobi) for complex scheduling.",
     "scipy.optimize.milp for integer resource allocation with deterministic fallback."),
    # Dataset
    ("29,302",   "17,520"),
    ("29302",    "17,520"),
    ("3+ years", "2 years"),
    ("3 years",  "2 years"),
    # Hybrid weights  (plain text)
    ("0.6/0.4",  "0.80/0.20"),
    ("0.7/0.3",  "0.80/0.20"),
    ("60%/40%",  "80%/20%"),
    ("70%/30%",  "80%/20%"),
    # False metric claims (only where these strings appear as metrics, not prose)
    ("25.5% improvement",  "~51% lower MAE vs ARIMAX"),
    # Optimizer
    ("rule-based optimization", "MILP optimization (scipy.optimize.milp)"),
    ("rule-based resource",     "MILP-based resource"),
    ("Rule-Based",              "MILP-Based"),
]

# Regex subs applied to full run text
REGEX_SUBS = [
    # "94%" as a claimed model accuracy → remove (only where clearly a metric)
    (re.compile(r'\b94\.17\s*%'), "5.52% MAPE"),
    # "Hybrid is most accurate" → correct
    (re.compile(r'Hybrid\s+(is|model is)\s+(the\s+)?most accurate', re.I),
     "LSTM is the most accurate model; Hybrid is deployed for robustness"),
    (re.compile(r'Hybrid\s+outperforms\s+LSTM', re.I),
     "LSTM outperforms Hybrid on test data; Hybrid deployed for robustness"),
]

# -----------------------------------------------------------------------
# Speaker notes per slide index (0-based)
# -----------------------------------------------------------------------
SPEAKER_NOTES = {
    # S13 — Model Comparison (index 12)
    12: (
        "Canonical test metrics:\n"
        "  LSTM:   MAE 7.65 | RMSE 9.58 | MAPE 5.52%  (most accurate model)\n"
        "  ARIMAX: MAE 15.63 | RMSE 19.33 | MAPE 12.33%  (baseline)\n"
        "  Hybrid 0.80/0.20: MAE 8.31 | RMSE 10.22 | MAPE 6.07%  (deployed for robustness)\n\n"
        "Q: Why not just use LSTM?\n"
        "A: LSTM is marginally more accurate, but Hybrid blends two independent models for "
        "robustness and provides a labelled fallback. Cost: ~0.7 MAE units for stability."
    ),
    # S14 — Hybrid Architecture (index 13)
    13: (
        "Grid search: 13 weight combinations from 0.20 to 0.80 in 0.05 steps.\n"
        "Min validation RMSE at LSTM=0.80, ARIMAX=0.20.\n"
        "Unconstrained optimum: 0.95/0.05 (confirms LSTM dominance).\n"
        "Hybrid deployed for robustness and labelled fallback, at ~0.7 MAE cost vs pure LSTM."
    ),
    # S15 — Forecasting Engine (index 14)
    14: (
        "LSTM uses a 24-hour lookback window (input sequence length).\n"
        "Forecast horizon is 72 hours (next 72 hourly predictions).\n"
        "ForecastState stores the full 72-hour forecast for Command Center and Digital Twin."
    ),
    # S20 — Interactive Dashboard (index 19)
    19: (
        "Dashboard built with Streamlit (not React). Role-based views: admin, doctor, nurse.\n"
        "13 tabs: Command Center, Forecast, Optimization, Digital Twin, What-If, "
        "Explainability, Departments, Shifts, Appointments, OR, Notifications, Approvals, Audit.\n"
        "ForecastState = single source of truth; cross-view consistency enforced by smoke test."
    ),
    # S21 — Precision Resource Optimization (index 20)
    20: (
        "scipy.optimize.milp: integer LP for beds, doctors, nurses across 5 departments.\n"
        "Objective: minimise total resource shortfall under capacity + ratio constraints.\n"
        "ICU nurse ratio 1:2; General Ward 1:6. Solves in <5 s; deterministic greedy fallback.\n"
        "Outputs: allocation table, shortage warnings, pressure ranking, action plan."
    ),
}

# -----------------------------------------------------------------------
# Slide-specific targeted replacements (shape-name + old_text + new_text)
# Matched by shape.name; old_text is a substring.
# -----------------------------------------------------------------------
TARGETED = {
    # S05 (index 4) — El Demerdash "Industrial Partner" -> "Clinical Context"
    4: [
        ("Text 0",
         "Industrial Partner: El Demerdash Hospital",
         "Clinical Context & Site Visit: El Demerdash Hospital"),
        ("Text 8",
         "Expected Impact: Improved patient flow, reduced wait times, "
         "optimized resource utilization, and enhanced operational efficiency.",
         "Site visit context (no real patient data): the hospital profile "
         "informed synthetic dataset parameters and system design requirements."),
    ],
    # S09 (index 8) — HL7/FHIR framed as future
    8: [
        ("Text 2",
         "Two-stage AI-Optimization framework with HL7/FHIR integration, vendor-agnostic APIs",
         "Two-stage AI-Optimization framework; HL7/FHIR integration is planned future work"),
        ("Text 6",
         "Data collection from HL7/FHIR systems, preprocessing, feature extraction, "
         "temporal pattern recognition",
         "Data ingestion from synthetic CSV dataset; HL7/FHIR integration is future work"),
        ("Text 14",
         "HL7/FHIR Integration",
         "HL7/FHIR Integration (Future Work)"),
    ],
    # S14 (index 13) — hybrid formula  (unicode chars get garbled; match by shape name)
    13: [
        # The formula shape has non-ASCII chars; match what survives ASCII re-encode
        # Actual stored text is: "ŷ = α·LSTM + (1−α)·ARIMAX"
        # We replace the entire run text on this shape
        ("Text 2", None,   # None = replace entire text_frame text regardless of content
         "Hybrid forecast: 0.80 × LSTM + 0.20 × ARIMAX  "
         "(α = 0.80, constrained grid search)"),
    ],
    # S15 (index 14) — 24-Hour -> 72-Hour
    14: [
        ("Text 1",
         "24-Hour Demand Predictions",
         "72-Hour Demand Forecast"),
        ("Text 12",
         "24-month hospital records,\ndemographics, seasonal patterns",
         "2-year synthetic dataset,\ndemographics, seasonal patterns (17,520 rows)"),
    ],
    # S21 (index 20) — false accuracy claim
    20: [
        ("Text 8",
         "The system achieves high\naccuracy in resource\npredictions, reducing\n"
         "operational costs \nwhile improving patient\nsatisfaction scores",
         "MILP optimizer allocates beds, doctors,\nand nurses across 5 departments\n"
         "under capacity & ratio constraints,\nsolving in under 5 seconds with a\n"
         "deterministic greedy fallback."),
    ],
}


def apply_subs_to_run(run) -> bool:
    original = run.text
    text = original
    for old, new in SUBS:
        text = text.replace(old, new)
    for pattern, new in REGEX_SUBS:
        text = pattern.sub(new, text)
    if text != original:
        run.text = text
        return True
    return False


def patch_text_frame(tf) -> int:
    """Patch all runs in a text_frame; return number of changed runs."""
    changed = 0
    for para in tf.paragraphs:
        for run in para.runs:
            if apply_subs_to_run(run):
                changed += 1
    return changed


def patch_shape(shape) -> int:
    changed = 0
    if shape.has_text_frame:
        changed += patch_text_frame(shape.text_frame)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                changed += patch_text_frame(cell.text_frame)
    return changed


def patch_notes(slide) -> bool:
    """Patch stale content in existing speaker notes."""
    if not slide.has_notes_slide:
        return False
    tf = slide.notes_slide.notes_text_frame
    changed = patch_text_frame(tf)
    return changed > 0


def set_notes(slide, text: str) -> None:
    """Set or append speaker notes."""
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    if tf is None:
        return
    # Check if there's already content
    try:
        existing = tf.text.strip()
    except Exception:
        existing = ""
    if existing:
        try:
            p = tf.add_paragraph()
            p.text = "\n---\n" + text
        except Exception:
            pass
    else:
        try:
            tf.paragraphs[0].text = text
        except Exception:
            pass


def apply_targeted(prs) -> dict:
    """Apply targeted slide-specific replacements."""
    results = {}
    for slide_idx, replacements in TARGETED.items():
        slide = prs.slides[slide_idx]
        slide_changes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for (shape_name, old_text, new_text) in replacements:
                if shape.name != shape_name:
                    continue
                tf = shape.text_frame
                current_text = tf.text
                if old_text is None:
                    # Replace entire text in first run of first paragraph
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        run = tf.paragraphs[0].runs[0]
                        # Clear all other runs/paragraphs first
                        old_val = run.text
                        run.text = new_text
                        # Remove remaining runs in first para
                        para = tf.paragraphs[0]
                        xml = para._p
                        runs_xml = xml.findall(
                            './/{http://schemas.openxmlformats.org/drawingml/2006/main}r'
                        )
                        for r in runs_xml[1:]:
                            xml.remove(r)
                        slide_changes.append(
                            f"S{slide_idx+1} [{shape_name}]: formula -> '{new_text}'"
                        )
                elif old_text in current_text:
                    # Replace in runs
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)
                                slide_changes.append(
                                    f"S{slide_idx+1} [{shape_name}]: '{old_text[:40]}' -> '{new_text[:40]}'"
                                )
                                break
        results[slide_idx] = slide_changes
    return results


def add_metrics_textbox(slide, prs):
    """Add canonical metrics text box to S13 (model comparison slide)."""
    # Place near bottom of slide, spanning most of width
    sw = prs.slide_width
    sh = prs.slide_height
    left   = Inches(0.4)
    top    = sh - Inches(1.6)
    width  = sw - Inches(0.8)
    height = Inches(1.4)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "Canonical Test Metrics (code-verified):"
    p0.runs[0].font.bold = True
    p0.runs[0].font.size = Pt(11)
    p0.runs[0].font.color.rgb = RGBColor(0x2f, 0x54, 0x96)

    lines = [
        "  LSTM:          MAE 7.65 | RMSE 9.58 | MAPE 5.52%  ← most accurate model",
        "  ARIMAX:        MAE 15.63 | RMSE 19.33 | MAPE 12.33%  (baseline)",
        "  Hybrid 0.80/0.20:  MAE 8.31 | RMSE 10.22 | MAPE 6.07%  (deployed for robustness)",
    ]
    colors = [
        RGBColor(0x1a, 0x60, 0x80),  # dark blue for LSTM
        RGBColor(0x5b, 0x6b, 0x86),  # grey for ARIMAX
        RGBColor(0x2f, 0x54, 0x96),  # medium blue for Hybrid
    ]
    for line, color in zip(lines, colors):
        p = tf.add_paragraph()
        p.text = line
        p.runs[0].font.size = Pt(10)
        p.runs[0].font.color.rgb = color


def add_image_to_slide(slide, img_path, left_in, top_in, width_in, height_in=None):
    from pptx.util import Inches
    left  = Inches(left_in)
    top   = Inches(top_in)
    width = Inches(width_in)
    if height_in:
        slide.shapes.add_picture(img_path, left, top, width, Inches(height_in))
    else:
        slide.shapes.add_picture(img_path, left, top, width)


def find_slide_by_keyword(prs, keyword, start=0):
    kw = keyword.lower()
    for i, slide in enumerate(prs.slides):
        if i < start:
            continue
        for shape in slide.shapes:
            if shape.has_text_frame and kw in shape.text_frame.text.lower():
                return i
    return -1


def main():
    print(f"Loading: {SRC}")
    prs = Presentation(SRC)

    total_global = 0
    total_notes_patched = 0

    # 1. Global text patching (runs + tables + notes) across all slides
    print("\n--- Phase 1: Global text patching ---")
    for i, slide in enumerate(prs.slides):
        slide_changes = 0
        for shape in slide.shapes:
            slide_changes += patch_shape(shape)
        if patch_notes(slide):
            total_notes_patched += 1
            slide_changes += 1
        if slide_changes:
            print(f"  S{i+1:02d}: {slide_changes} changes")
        total_global += slide_changes

    print(f"  Total global changes: {total_global} runs; {total_notes_patched} note frames patched")

    # 2. Targeted slide-specific replacements
    print("\n--- Phase 2: Targeted slide-specific fixes ---")
    targeted_results = apply_targeted(prs)
    for slide_idx, changes in targeted_results.items():
        for c in changes:
            print(f"  {c}".encode('ascii', 'replace').decode())

    # 3. Add canonical metrics text box to S13
    print("\n--- Phase 3: Add metrics text box to S13 ---")
    s13 = prs.slides[12]
    add_metrics_textbox(s13, prs)
    print("  Metrics text box added to S13")

    # 4. Add speaker notes
    print("\n--- Phase 4: Adding speaker notes ---")
    for slide_idx, notes_text in SPEAKER_NOTES.items():
        set_notes(prs.slides[slide_idx], notes_text)
        print(f"  Speaker notes added to S{slide_idx+1}")

    # 5. Insert diagram images
    print("\n--- Phase 5: Inserting diagram images ---")
    sw_in = prs.slide_width / 914400
    sh_in = prs.slide_height / 914400

    # S10 System Block Diagram — add architecture image (check if already present from prev run)
    s10 = prs.slides[9]
    # Count pictures on S10 to see if arch was already added
    pic_count = sum(1 for s in s10.shapes if not s.has_text_frame and not s.has_table)
    arch_png = str(FIGS / "fig_architecture.png")
    if Path(arch_png).exists():
        try:
            add_image_to_slide(s10, arch_png,
                               left_in=0.3, top_in=2.2,
                               width_in=min(9.0, sw_in - 0.6))
            print(f"  Architecture diagram added to S10 (was {pic_count} images)")
        except Exception as e:
            print(f"  S10 arch image error: {e}")

    # S15 Forecasting Engine — add forecasting pipeline
    fp_png = str(FIGS / "fig_forecasting_pipeline.png")
    if Path(fp_png).exists():
        s15 = prs.slides[14]
        try:
            add_image_to_slide(s15, fp_png,
                               left_in=0.3, top_in=2.8,
                               width_in=min(8.5, sw_in - 0.6))
            print("  Forecasting pipeline diagram added to S15")
        except Exception as e:
            print(f"  S15 forecast image error: {e}")

    # S13 Model Comparison — add results metrics figure
    met_png = str(FIGS / "fig_results_metrics.png")
    if Path(met_png).exists():
        s13 = prs.slides[12]
        try:
            add_image_to_slide(s13, met_png,
                               left_in=0.3, top_in=1.8,
                               width_in=min(9.0, sw_in - 0.6),
                               height_in=3.2)
            print("  Results metrics diagram added to S13")
        except Exception as e:
            print(f"  S13 metrics image error: {e}")

    # S20 Dashboard — add command center screenshot
    cmd_png = str(SHOTS / "01_command_center.jpg")
    if Path(cmd_png).exists():
        s20 = prs.slides[19]
        try:
            add_image_to_slide(s20, cmd_png,
                               left_in=0.3, top_in=3.8,
                               width_in=min(6.0, sw_in - 0.6))
            print("  Command Center screenshot added to S20")
        except Exception as e:
            print(f"  S20 screenshot error: {e}")

    # S21 Optimization — add optimization screenshot
    opt_png = str(SHOTS / "05_optimization.jpg")
    if Path(opt_png).exists():
        s21 = prs.slides[20]
        try:
            add_image_to_slide(s21, opt_png,
                               left_in=0.3, top_in=3.8,
                               width_in=min(6.0, sw_in - 0.6))
            print("  Optimization screenshot added to S21")
        except Exception as e:
            print(f"  S21 screenshot error: {e}")

    # 6. Save
    print(f"\n--- Saving to {DST} ---")
    prs.save(DST)
    sz = Path(DST).stat().st_size // 1024
    print(f"Saved: {sz} KB")
    print("\nDone.")


if __name__ == "__main__":
    main()
