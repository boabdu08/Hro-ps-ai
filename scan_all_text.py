"""Scan ALL text in every slide shape (including speaker notes) for stale metrics."""
from pptx import Presentation
import re

prs = Presentation(r'D:\Hospital Resource Optimization with AI Final Hro.pptx')

# Patterns to flag
FLAG_PATTERNS = [
    r'94\.17', r'\b94%', r'29,?302', r'25\.5%', r'18.22%', r'3\+?\s*years',
    r'168.?[Hh]our', r'rule.?based', r'OR.?Tools', r'CP.?SAT', r'TimescaleDB',
    r'React', r'HL7/FHIR(?! integration \(future)',
    r'0\.6/0\.4', r'0\.7/0\.3', r'60%/40%', r'70%/30%',
    r'7\.32', r'9\.83', r'7\.76', r'\b6\.6\b', r'\b8\.1\b', r'\b4\.9%',
    r'24.Hour Demand', r'24-Hour Demand',
    r'productivity', r'improvement',
    r'\balpha\s*=\s*0\.[67]', r'alpha.*0\.6', r'a.*=.*0\.7',
    r'most accurate.*[Hh]ybrid', r'[Hh]ybrid.*most accurate',
]
combined = re.compile('|'.join(FLAG_PATTERNS), re.I)

print('=== STALE TEXT SCAN ACROSS ALL SLIDES ===\n')

for i, slide in enumerate(prs.slides, 1):
    flagged = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text
                    if combined.search(t):
                        flagged.append(f'  [shape={shape.name}] "{t.strip()[:80]}"')
        if shape.has_table:
            for ri, row in enumerate(shape.table.rows):
                for ci, cell in enumerate(row.cells):
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            t = run.text
                            if combined.search(t):
                                flagged.append(f'  [TABLE row={ri} col={ci}] "{t.strip()[:80]}"')
    # Check speaker notes
    if slide.has_notes_slide:
        notes_tf = slide.notes_slide.notes_text_frame
        for para in notes_tf.paragraphs:
            for run in para.runs:
                t = run.text
                if combined.search(t):
                    flagged.append(f'  [NOTES] "{t.strip()[:80]}"')

    if flagged:
        print(f'S{i:02d}:')
        for f in flagged:
            print(f)
        print()

print('=== DONE ===')
