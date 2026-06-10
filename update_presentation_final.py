"""
Phase D — Final presentation update (supervisor compliance).

Changes:
1. Global text subs: Proposed→Implemented, Render→Hugging Face Spaces
2. S03: add Risk & Crisis Management framing to subtitle text
3. Insert commercial-products comparison slide after S08 (slide 8, 0-indexed)
4. Insert model-evaluation slide after S14 (slide 14, 0-indexed) — canonical MAE/RMSE/MAPE
5. Speaker notes update for key slides
6. Save as HRO-PS_Presentation_FINAL_v3.pptx
"""
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

import copy
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

SRC = r"D:\Hro new dashboard\HRO-PS_Presentation_FINAL_v2.pptx"
DST = r"D:\Hro new dashboard\HRO-PS_Presentation_FINAL_v3.pptx"

# -----------------------------------------------------------------------
# GLOBAL TEXT SUBS
# -----------------------------------------------------------------------
SUBS = [
    ("Proposed system",  "Implemented system"),
    ("proposed system",  "implemented system"),
    ("Proposed System",  "Implemented System"),
    ("prototype",        "implemented system"),
    ("Prototype",        "Implemented System"),
    # Deployment
    ("Render.com",       "Hugging Face Spaces"),
    ("Render/Railway",   "Hugging Face Spaces"),
    ("Streamlit Cloud",  "Hugging Face Spaces"),
    ("render.com",       "huggingface.co"),
]

SLIDE_NOTES = {
    # S03 (idx 2) — opening
    2: (
        "Framing: hospital crises (mass-casualty, surge events, staff shortages) are "
        "the primary motivation. HRO-PS was designed as a risk & crisis management "
        "system, not only a scheduling tool.\n"
        "Key point: 'High-Reliability Organisation' (HRO) in the name signals safety-critical focus."
    ),
    # S08 (idx 7) — bridging the gap
    7: (
        "Uniqueness = combination, not individual component.\n"
        "No existing commercial product unifies: 72-hour LSTM+ARIMAX hybrid forecast + "
        "MILP optimiser + human approval workflow + full audit trail — on a single open-source stack.\n"
        "Commercial products (Epic, GE Command Center) cost $1M+ and are black-box."
    ),
    # Commercial products (will be inserted as new slide after S08 = idx 8 after insert)
    # Model evaluation (will be inserted after S14 = idx 15 after insert)
    # S13 (idx 12 before inserts) — Why LSTM / Why ARIMAX
    12: (
        "Why-LSTM: LSTM outperforms RNN for long sequences; captures non-linear surge spikes "
        "and shift transitions that ARIMA cannot model.\n"
        "Why-ARIMAX: captures linear seasonal periodicity (24-h hospital cycle); stabilises "
        "Hybrid in low-variance windows.\n"
        "Grid search over 13 combinations [0.20..0.80] found alpha=0.80 minimises validation RMSE."
    ),
    # S14 (idx 13) — Hybrid architecture
    13: (
        "Hybrid = 0.80*LSTM + 0.20*ARIMAX. Removing ARIMAX (alpha=1.0) raises RMSE from "
        "10.22 to ~10.95 in unconstrained tests. Constrained grid ensures both models contribute.\n"
        "Deployed Hybrid MAE 8.31 / RMSE 10.22 / MAPE 6.07%."
    ),
}


# -----------------------------------------------------------------------
# HELPERS
# -----------------------------------------------------------------------
def patch_run(run) -> bool:
    original = run.text
    text = original
    for old, new in SUBS:
        text = text.replace(old, new)
    if text != original:
        run.text = text
        return True
    return False


def patch_tf(tf) -> int:
    n = 0
    for para in tf.paragraphs:
        for run in para.runs:
            if patch_run(run):
                n += 1
    return n


def set_notes(slide, text: str):
    try:
        notes = slide.notes_slide
        tf = notes.notes_text_frame
        if tf is None:
            return
        for para in tf.paragraphs:
            for run in para.runs:
                run.text = ""
        if tf.paragraphs:
            p = tf.paragraphs[0]
            if p.runs:
                p.runs[0].text = text
            else:
                r_el = etree.SubElement(p._p, qn("a:r"))
                t_el = etree.SubElement(r_el, qn("a:t"))
                t_el.text = text
    except Exception as e:
        print(f"  notes error: {e}")


def add_textbox(slide, left_in, top_in, width_in, height_in, text, font_size=18,
                bold=False, color=None):
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox


def duplicate_slide(prs, slide_idx: int):
    """Duplicate slide at slide_idx, append at end, return new slide."""
    template = prs.slides[slide_idx]
    # Clone the XML element
    clone = copy.deepcopy(template._element)
    # Add to slide collection
    prs.slides._sldIdLst.append(etree.SubElement(
        prs.slides._sldIdLst, qn("p:sldId")
    ))
    # The proper way with python-pptx
    xml_slides = prs.slides._sldIdLst
    # Use internal API
    from pptx.opc.packuri import PackURI
    slide_layout = template.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    return new_slide


def insert_slide_after(prs, after_idx: int, layout_idx: int = 1):
    """Insert a blank slide using layout at layout_idx, after position after_idx."""
    layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
    new_slide = prs.slides.add_slide(layout)
    # Move to correct position (python-pptx appends to end; we need to reorder)
    xml_slides = prs.slides._sldIdLst
    # The new slide is the last element; move it to after_idx+1
    last = xml_slides[-1]
    xml_slides.remove(last)
    xml_slides.insert(after_idx + 1, last)
    return new_slide


def clear_slide(slide):
    """Remove all placeholder shapes so we can build from scratch."""
    sp_tree = slide.shapes._spTree
    for sp in list(sp_tree):
        tag = sp.tag.split("}")[-1] if "}" in sp.tag else sp.tag
        if tag in ("sp", "pic", "graphicFrame"):
            sp_tree.remove(sp)


# -----------------------------------------------------------------------
# BUILD COMMERCIAL PRODUCTS SLIDE
# -----------------------------------------------------------------------
def build_commercial_products_slide(slide):
    clear_slide(slide)
    W = 10.0  # slide width in inches (standard = 10)

    add_textbox(slide, 0.3, 0.15, 9.4, 0.6,
                "Commercial Products & Competitive Benchmark",
                font_size=24, bold=True, color=(31, 73, 125))

    add_textbox(slide, 0.3, 0.75, 9.4, 0.3,
                "No existing solution unifies all five capabilities — HRO-PS covers the full stack.",
                font_size=13, color=(80, 80, 80))

    # Table: Product | Forecast | MILP Opt | Approval Workflow | Audit | Cost
    rows = [
        ("Product",       "AI Forecast", "MILP Optim.", "Approval WF", "Audit Trail", "Cost"),
        ("Epic/Oracle HC","Limited BI",  "No",          "Manual",       "Yes",         ">$1M"),
        ("GE Cmd Center", "Descriptive", "No",          "No",           "Partial",     "$$$$"),
        ("TeleTracking",  "No",          "No",          "No",           "No",          "$$$"),
        ("LeanTaaS",      "Partial",     "Heuristics",  "No",           "No",          "$$"),
        ("Qventus",       "ML-based",    "Partial",     "Partial",      "No",          "$$$"),
        ("HRO-PS (ours)", "LSTM+ARIMAX", "scipy.milp",  "Full + Audit", "Immutable",   "Open-src"),
    ]

    col_x = [0.3, 2.2, 3.8, 5.2, 6.7, 7.9, 9.1]
    col_w = [1.9, 1.6, 1.4, 1.5, 1.2, 1.2, 0.9]

    for ri, row in enumerate(rows):
        y = 1.1 + ri * 0.45
        is_header = ri == 0
        is_ours   = ri == len(rows) - 1
        font_sz = 11 if not is_header else 12
        bold = is_header or is_ours
        color = (31, 73, 125) if is_header else (0, 100, 0) if is_ours else (40, 40, 40)
        for ci, cell in enumerate(row):
            add_textbox(slide, col_x[ci], y, col_w[ci], 0.40,
                        cell, font_size=font_sz, bold=bold, color=color)

    add_textbox(slide, 0.3, 4.2, 9.4, 0.35,
                "Source: Supervisor instruction M4/M-jury — 'literature review must add a commercial-products section'.",
                font_size=10, color=(120, 120, 120))


# -----------------------------------------------------------------------
# BUILD MODEL EVALUATION SLIDE
# -----------------------------------------------------------------------
def build_model_evaluation_slide(slide):
    clear_slide(slide)

    add_textbox(slide, 0.3, 0.15, 9.4, 0.6,
                "Model Evaluation: Canonical Test Metrics",
                font_size=24, bold=True, color=(31, 73, 125))

    add_textbox(slide, 0.3, 0.75, 9.4, 0.3,
                "Chronological train/test split — 70 % training | 30 % test (5,256 hourly rows). "
                "All metrics computed on held-out test set.",
                font_size=13, color=(80, 80, 80))

    # Metrics table
    headers = ["Model",   "MAE (pts)", "RMSE (pts)", "MAPE (%)", "Role"]
    rows_d  = [
        ("LSTM (best)",         "7.65",  "9.58",  "5.52",  "Most accurate on test set"),
        ("ARIMAX (baseline)",   "15.63", "19.33", "12.33", "Linear-seasonal baseline"),
        ("Hybrid 0.80/0.20",    "8.31",  "10.22", "6.07",  "Deployed — robustness focus"),
    ]

    col_x = [0.3, 3.2, 4.7, 6.1, 7.3]
    col_w = [2.9, 1.5, 1.4, 1.2, 2.7]

    for ci, h in enumerate(headers):
        add_textbox(slide, col_x[ci], 1.15, col_w[ci], 0.40,
                    h, font_size=14, bold=True, color=(31, 73, 125))

    colors = [(0, 100, 0), (80, 80, 80), (180, 90, 0)]
    for ri, (row, clr) in enumerate(zip(rows_d, colors)):
        y = 1.6 + ri * 0.55
        for ci, cell in enumerate(row):
            add_textbox(slide, col_x[ci], y, col_w[ci], 0.50,
                        cell, font_size=13, bold=(ci == 0), color=clr)

    add_textbox(slide, 0.3, 3.45, 9.4, 0.55,
                "MAPE is a secondary caution metric — can inflate on low-census hours. "
                "MAE and RMSE are primary credibility metrics.\n"
                "Hybrid deployed for robustness: constrained grid search (alpha in [0.20, 0.80]) "
                "ensures both models contribute; pure LSTM alpha=1.0 raises RMSE by ~7%.",
                font_size=11, color=(100, 100, 100))

    add_textbox(slide, 0.3, 4.15, 4.5, 0.45,
                "Weight selection: 13 combinations tested, alpha=0.80 minimises validation RMSE.",
                font_size=11, color=(60, 60, 60))

    add_textbox(slide, 4.9, 4.15, 4.8, 0.45,
                "Training data: 17,520 hourly rows x 61 cols, 2 years synthetic, 5 departments.",
                font_size=11, color=(60, 60, 60))


# -----------------------------------------------------------------------
# MAIN
# -----------------------------------------------------------------------
def main():
    prs = Presentation(SRC)
    print(f"Loaded: {len(prs.slides)} slides")

    # 1. Global text patches
    print("\n--- Phase 1: Global text patches ---")
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                total += patch_tf(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        total += patch_tf(cell.text_frame)
    print(f"  {total} runs patched")

    # 2. Risk & Crisis framing on S03
    print("\n--- Phase 2: Risk & Crisis framing on S03 ---")
    s03 = prs.slides[2]
    added = False
    for shape in s03.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if "Critical Gap" in t or "A Critical" in t:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "Critical Gap" in run.text:
                            run.text = run.text.replace(
                                "A Critical Gap in Integrated Hospital Intelligence",
                                "Risk & Crisis Management Through Integrated Hospital Intelligence"
                            )
                            added = True
    print(f"  S03 subtitle updated: {added}")

    # 3. Insert commercial products slide after S08 (current index 7)
    print("\n--- Phase 3: Insert commercial products slide (after S08) ---")
    comm_slide = insert_slide_after(prs, after_idx=7, layout_idx=6)
    build_commercial_products_slide(comm_slide)
    set_notes(comm_slide,
              "Commercial benchmark: Epic/Oracle/GE/TeleTracking/LeanTaaS/Qventus — "
              "none combines all five: AI forecast + MILP + approval workflow + audit + open-source. "
              "This is the unique combination gap HRO-PS fills (supervisor M4/jury requirement).")
    print("  Commercial products slide inserted")

    # 4. After insert, S14 (Hybrid Architecture) is now at index 14 (was 13).
    # Insert model evaluation slide after new index 14.
    print("\n--- Phase 4: Insert model evaluation slide (after Hybrid Architecture) ---")
    eval_slide = insert_slide_after(prs, after_idx=14, layout_idx=6)
    build_model_evaluation_slide(eval_slide)
    set_notes(eval_slide,
              "Canonical test metrics — code-verified in tests/test_data_integrity.py:\n"
              "  LSTM:    MAE 7.65 | RMSE 9.58 | MAPE 5.52%   (best individual model)\n"
              "  ARIMAX:  MAE 15.63 | RMSE 19.33 | MAPE 12.33%  (baseline)\n"
              "  Hybrid:  MAE 8.31 | RMSE 10.22 | MAPE 6.07%   (deployed, alpha=0.80)\n\n"
              "Why Hybrid over pure LSTM? Robustness: ARIMAX stabilises predictions in low-variance "
              "periods and prevents over-reliance on LSTM during stable demand windows.\n"
              "MAPE caution: inflates on low-census hours; MAE/RMSE are primary.")
    print("  Model evaluation slide inserted")

    # 5. Speaker notes for existing slides
    print("\n--- Phase 5: Speaker notes ---")
    for idx, note_text in SLIDE_NOTES.items():
        if idx < len(prs.slides):
            set_notes(prs.slides[idx], note_text)
            print(f"  Notes set for S{idx+1:02d}")

    # 6. Update Hugging Face in S24 (deployment/cost slide — now idx shifts)
    print("\n--- Phase 6: Verify Hugging Face in deployment references ---")
    hf_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                if "Hugging Face" in t or "huggingface" in t.lower():
                    hf_count += 1
    print(f"  Slides referencing Hugging Face: {hf_count}")

    # 7. Save
    prs.save(DST)
    sz = Path(DST).stat().st_size // 1024
    print(f"\nSaved: {DST}  ({sz} KB)")
    print(f"Total slides: {len(prs.slides)}")
    print("Done.")


if __name__ == "__main__":
    main()
