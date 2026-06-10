"""Get full text of specific shapes on key slides."""
from pptx import Presentation

prs = Presentation(r'D:\Hospital Resource Optimization with AI Final Hro.pptx')

checks = {
    4:  None,    # S05 - El Demerdash
    12: None,    # S13 - model comparison (all text)
    13: None,    # S14 - hybrid formula
    20: None,    # S21 - optimization / accuracy claims
}

for slide_idx in checks:
    slide = prs.slides[slide_idx]
    print(f'=== S{slide_idx+1} ALL TEXT ===')
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                print(f'  [{shape.name}]: {repr(t[:200])}'.encode('ascii', 'replace').decode())
        if shape.has_table:
            for ri, row in enumerate(shape.table.rows):
                for ci, cell in enumerate(row.cells):
                    ct = cell.text_frame.text.strip()
                    if ct:
                        print(f'  [TABLE r{ri}c{ci}]: {repr(ct[:100])}'.encode('ascii','replace').decode())
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            print(f'  [NOTES]: {repr(notes[:200])}'.encode('ascii','replace').decode())
    print()
