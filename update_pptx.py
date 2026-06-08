"""
Apply corrected content (from HRO-PS_Poster_Content.md and
HRO-PS_Presentation_Content_and_Sync.md) to the two source PPTX files.

Rules:
- Preserve ALL existing theme, fonts, colours, slide layout, and master.
- Only replace text runs that contain stale strings; never touch runs that
  look fine.
- Insert diagram PNGs where figure slots exist (System Block Diagram slide,
  Forecasting slide, Results slide).
- Save as NEW files (never overwrite originals).
"""
import re
import sys
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

BASE = Path(r"D:\Hro new dashboard")
FIGS = BASE / "thesis_figures"
SHOTS = FIGS / "screenshots"

# ------------------------------------------------------------------
# Canonical replacements — (old_pattern, new_text) pairs
# Applied to every text run in every slide.
# old_pattern is a substring match (case-sensitive except where noted).
# ------------------------------------------------------------------
TEXT_SUBS = [
    # Dataset corrections
    ("29,302", "17,520"),
    ("29302",  "17,520"),
    ("3+ years", "2 years"),
    ("3 years",  "2 years"),
    ("168-hour",  "24-hour"),
    ("168 hour",  "24-hour"),
    ("168-Hour",  "24-Hour"),
    ("168 Hour",  "24-Hour"),
    ("168h",      "24h"),

    # Hybrid weight corrections
    ("0.6/0.4",  "0.80/0.20"),
    ("0.7/0.3",  "0.80/0.20"),
    ("60%/40%",  "80%/20%"),
    ("70%/30%",  "80%/20%"),
    ("α = 0.6",  "α = 0.80"),
    ("α = 0.7",  "α = 0.80"),
    ("α=0.6",    "α = 0.80"),
    ("α=0.7",    "α = 0.80"),

    # Accuracy / metric corrections  (remove false claims)
    ("94.17%",   "5.52% MAPE"),
    ("94%",      "LSTM best model"),
    ("25.5% improvement",  "~51% lower MAE vs ARIMAX"),
    ("18-22% productivity", "~51% lower MAE vs ARIMAX"),
    ("18–22% productivity", "~51% lower MAE vs ARIMAX"),

    # Optimizer correction
    ("rule-based optimization", "MILP optimization (scipy.optimize.milp)"),
    ("rule-based resource",     "MILP-based resource"),
    ("Rule-Based",              "MILP-Based"),
    ("rule-based",              "MILP-based"),
    ("OR-Tools",                "scipy.optimize.milp"),
    ("CP-SAT",                  "scipy.optimize.milp"),

    # Model spec corrections
    ("TimescaleDB",  "PostgreSQL"),
    ("React frontend", "Streamlit dashboard"),

    # Test count
    ("108 unit tests", "128 unit tests"),
    ("108 tests",      "128 tests"),

    # Stack corrections
    ("HL7/FHIR integration", "HL7/FHIR integration (future work)"),
]

# Regex-level replacements (applied to the full paragraph text)
REGEX_SUBS = [
    # Hybrid "best/most accurate" → LSTM is best
    (re.compile(r"Hybrid\s+(is|model is)\s+(the\s+)?most accurate", re.I),
     "LSTM is the most accurate model; Hybrid is deployed for robustness"),
    (re.compile(r"Hybrid\s+outperforms\s+LSTM", re.I),
     "LSTM outperforms the Hybrid on test data; Hybrid deployed for robustness"),
    # Dataset row count
    (re.compile(r"\b29[,\s]?302\b"),   "17,520"),
    (re.compile(r"\b29302\b"),          "17,520"),
]

# Phrases to append a clarifying note to (sentence-level)
APPEND_NOTES = {
    "Hybrid (constrained blend)": " — LSTM 0.80 / ARIMAX 0.20",
}


def apply_text_subs(text: str) -> str:
    """Apply all simple and regex substitutions to a text string."""
    for old, new in TEXT_SUBS:
        text = text.replace(old, new)
    for pattern, new in REGEX_SUBS:
        text = pattern.sub(new, text)
    return text


def patch_run(run) -> bool:
    """Return True if text was changed."""
    original = run.text
    patched = apply_text_subs(original)
    if patched != original:
        run.text = patched
        return True
    return False


def patch_slide(slide) -> int:
    """Patch all text frames in a slide; return count of changed runs."""
    changes = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if patch_run(run):
                    changes += 1
    return changes


def add_image_to_slide(slide, img_path: str,
                        left_in: float, top_in: float,
                        width_in: float, height_in: float = None):
    """Add a PNG image to the slide without disturbing existing shapes."""
    from pptx.util import Inches
    left  = Inches(left_in)
    top   = Inches(top_in)
    width = Inches(width_in)
    if height_in:
        height = Inches(height_in)
        slide.shapes.add_picture(img_path, left, top, width, height)
    else:
        slide.shapes.add_picture(img_path, left, top, width)


def find_slide_by_keyword(prs, keyword: str, start: int = 0):
    """Return the first slide (index) whose text contains keyword."""
    kw = keyword.lower()
    for i, slide in enumerate(prs.slides):
        if i < start:
            continue
        for shape in slide.shapes:
            if shape.has_text_frame:
                if kw in shape.text_frame.text.lower():
                    return i
    return -1


# ------------------------------------------------------------------
# 1. Poster (poster msa.pptx)  →  HRO-PS_Poster_FINAL.pptx
# ------------------------------------------------------------------

def update_poster():
    src = r"D:\poster msa.pptx"
    dst = str(BASE / "HRO-PS_Poster_FINAL.pptx")
    prs = Presentation(src)

    total_changes = 0
    for slide in prs.slides:
        total_changes += patch_slide(slide)

    # Insert architecture diagram (top-left area, below header)
    slide = prs.slides[0]
    arch_png = str(FIGS / "fig_architecture.png")
    if Path(arch_png).exists():
        try:
            add_image_to_slide(slide, arch_png,
                               left_in=0.2, top_in=5.5,
                               width_in=4.5)
        except Exception as e:
            print(f"  Poster arch image: {e}")

    # Insert results metrics diagram
    metrics_png = str(FIGS / "fig_results_metrics.png")
    if Path(metrics_png).exists():
        try:
            add_image_to_slide(slide, metrics_png,
                               left_in=5.0, top_in=5.5,
                               width_in=4.5)
        except Exception as e:
            print(f"  Poster metrics image: {e}")

    prs.save(dst)
    sz = Path(dst).stat().st_size // 1024
    print(f"Poster saved: {dst}  ({sz} KB, {total_changes} text changes)")


# ------------------------------------------------------------------
# 2. Presentation (Hospital Resource Optimization...Final Hro.pptx)
#    →  HRO-PS_Presentation_FINAL.pptx
# ------------------------------------------------------------------

def update_presentation():
    src = r"D:\Hospital Resource Optimization with AI Final Hro.pptx"
    dst = str(BASE / "HRO-PS_Presentation_FINAL.pptx")
    prs = Presentation(src)

    total_changes = 0
    for slide in prs.slides:
        total_changes += patch_slide(slide)

    w = prs.slide_width
    h = prs.slide_height

    def slide_inches():
        return w / 914400, h / 914400  # EMU → inches

    sw, sh = slide_inches()

    # Insert architecture diagram on "System Block Diagram" slide
    idx = find_slide_by_keyword(prs, "system block diagram")
    if idx >= 0:
        arch_png = str(FIGS / "fig_architecture.png")
        if Path(arch_png).exists():
            try:
                add_image_to_slide(prs.slides[idx], arch_png,
                                   left_in=0.3, top_in=2.2,
                                   width_in=min(9.0, sw - 0.6))
                print(f"  Architecture diagram added to slide {idx+1}")
            except Exception as e:
                print(f"  Arch image S{idx+1}: {e}")

    # Insert forecasting pipeline diagram on "Hybrid Forecasting Engine" slide
    idx = find_slide_by_keyword(prs, "hybrid forecasting engine")
    if idx < 0:
        idx = find_slide_by_keyword(prs, "forecasting engine")
    if idx >= 0:
        fp_png = str(FIGS / "fig_forecasting_pipeline.png")
        if Path(fp_png).exists():
            try:
                add_image_to_slide(prs.slides[idx], fp_png,
                                   left_in=0.3, top_in=2.5,
                                   width_in=min(8.5, sw - 0.6))
                print(f"  Forecasting pipeline diagram added to slide {idx+1}")
            except Exception as e:
                print(f"  Forecast image S{idx+1}: {e}")

    # Insert results metrics diagram on "Algorithm Selection" or "Model Comparison" slide
    idx = find_slide_by_keyword(prs, "model comparison")
    if idx < 0:
        idx = find_slide_by_keyword(prs, "algorithm selection")
    if idx >= 0:
        met_png = str(FIGS / "fig_results_metrics.png")
        if Path(met_png).exists():
            try:
                add_image_to_slide(prs.slides[idx], met_png,
                                   left_in=0.3, top_in=2.5,
                                   width_in=min(8.5, sw - 0.6))
                print(f"  Results metrics diagram added to slide {idx+1}")
            except Exception as e:
                print(f"  Metrics image S{idx+1}: {e}")

    # Insert Optimization screenshot on "Precision Resource Optimization" slide
    idx = find_slide_by_keyword(prs, "precision resource optimization")
    if idx < 0:
        idx = find_slide_by_keyword(prs, "resource optimization")
    if idx >= 0:
        opt_png = str(SHOTS / "05_optimization.jpg")
        if Path(opt_png).exists():
            try:
                add_image_to_slide(prs.slides[idx], opt_png,
                                   left_in=0.3, top_in=3.5,
                                   width_in=min(6.0, sw - 0.6))
                print(f"  Optimization screenshot added to slide {idx+1}")
            except Exception as e:
                print(f"  Opt image S{idx+1}: {e}")

    # Insert Dashboard screenshot on "Interactive Operational Dashboard" slide
    idx = find_slide_by_keyword(prs, "operational dashboard")
    if idx >= 0:
        cmd_png = str(SHOTS / "01_command_center.jpg")
        if Path(cmd_png).exists():
            try:
                add_image_to_slide(prs.slides[idx], cmd_png,
                                   left_in=0.3, top_in=3.5,
                                   width_in=min(6.0, sw - 0.6))
                print(f"  Command Center screenshot added to slide {idx+1}")
            except Exception as e:
                print(f"  CMD image S{idx+1}: {e}")

    prs.save(dst)
    sz = Path(dst).stat().st_size // 1024
    print(f"Presentation saved: {dst}  ({sz} KB, {total_changes} text changes)")


if __name__ == "__main__":
    print("=== Updating poster ===")
    update_poster()
    print("\n=== Updating presentation ===")
    update_presentation()
    print("\nDone.")
